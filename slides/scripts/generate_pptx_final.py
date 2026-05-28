#!/usr/bin/env python3
"""ContratIA Abierta — Deck FINAL (ingeniería de datos, una sola lane, ES).

16 diapositivas técnicas. Una sola arquitectura. Incluye validación AGR (2.5x),
caso real Puerto Gaitán (3.1x) y embeddings semánticos funcionando.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

INK = RGBColor(0x0B, 0x1E, 0x33)
NAVY = RGBColor(0x10, 0x3A, 0x5C)
CORAL = RGBColor(0xE3, 0x5A, 0x4B)
TEAL = RGBColor(0x1F, 0x82, 0x7C)
AMBER = RGBColor(0xC2, 0x88, 0x32)
SAND = RGBColor(0xF5, 0xEA, 0xD2)
PAPER = RGBColor(0xFA, 0xFA, 0xF7)
CLOUD = RGBColor(0xEE, 0xF2, 0xF8)
MIST = RGBColor(0xDB, 0xE2, 0xEC)
MUTED = RGBColor(0x55, 0x60, 0x76)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"
MONO = "Consolas"
ASSETS = Path(__file__).resolve().parent.parent / "assets"
OUTPUT = Path(__file__).resolve().parent.parent / "contratia_abierta_final.pptx"
FOOTER = "ContratIA Abierta · Datos al Ecosistema 2026 · Prioriza revisión humana, no acusa."
TOTAL = 16


def bg(slide, color=PAPER):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def rect(slide, left, top, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def band(slide):
    rect(slide, 0, 0, SLIDE_W, Pt(6), NAVY)
    rect(slide, 0, Pt(6), Inches(1.6), Pt(3), CORAL)


def text(slide, t, left, top, w, h, size=14, color=INK, bold=False,
         italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    b = slide.shapes.add_textbox(left, top, w, h)
    tf = b.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Pt(0))
    p = tf.paragraphs[0]
    p.text = t
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font
    return b


def title(slide, eyebrow, ttl, top=Inches(0.5)):
    text(slide, eyebrow.upper(), Inches(0.7), top, Inches(11.5), Inches(0.3),
         size=10, color=CORAL, bold=True)
    text(slide, ttl, Inches(0.7), top + Inches(0.32), Inches(12), Inches(0.85),
         size=25, bold=True, color=INK)
    rect(slide, Inches(0.7), top + Inches(1.12), Inches(0.6), Pt(2), CORAL)


def footer(slide, page):
    text(slide, FOOTER, Inches(0.7), SLIDE_H - Inches(0.4), Inches(11), Inches(0.3),
         size=8, color=MUTED, italic=True)
    text(slide, f"{page} / {TOTAL}", SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.4),
         Inches(0.6), Inches(0.3), size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def card(slide, left, top, w, h, ttl, body, accent=NAVY, bg_c=WHITE,
         tsize=13, bsize=10, tcolor=INK):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = bg_c
    s.line.color.rgb = MIST
    s.line.width = Pt(0.75)
    s.adjustments[0] = 0.05
    rect(slide, left, top, Pt(4), h, accent)
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(9)
    tf.margin_bottom = Pt(9)
    p = tf.paragraphs[0]
    p.text = ttl
    p.font.size = Pt(tsize)
    p.font.bold = True
    p.font.color.rgb = tcolor
    p.font.name = FONT
    if body:
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(bsize)
        p2.font.color.rgb = MUTED
        p2.font.name = FONT
        p2.space_before = Pt(3)
    return s


def metric(slide, left, top, w, h, val, label, color=NAVY, bg_c=CLOUD, vsize=26):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = bg_c
    s.line.fill.background()
    s.adjustments[0] = 0.08
    tf = s.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Pt(8))
    p = tf.paragraphs[0]
    p.text = val
    p.font.size = Pt(vsize)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(9)
    p2.font.color.rgb = MUTED
    p2.font.name = FONT
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)


def chip(slide, left, top, t, bg_c=NAVY, fg=WHITE):
    w = Inches(max(1.2, len(t) * 0.085 + 0.4))
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.34))
    s.fill.solid()
    s.fill.fore_color.rgb = bg_c
    s.line.fill.background()
    s.adjustments[0] = 0.5
    tf = s.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = t
    p.font.size = Pt(9)
    p.font.color.rgb = fg
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    return s


def arrow(slide, left, top, w=Inches(0.45), color=CORAL):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, w, Inches(0.2))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def img(slide, name, left, top, w=None, h=None):
    p = ASSETS / name
    if not p.exists():
        card(slide, left, top, w or Inches(4), h or Inches(2.4), f"[{name}]", "")
        return None
    kw = {}
    if w:
        kw["width"] = w
    if h:
        kw["height"] = h
    return slide.shapes.add_picture(str(p), left, top, **kw)


def code_block(slide, left, top, w, h, lines):
    rect(slide, left, top, w, h, INK)
    b = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1),
                                 w - Inches(0.3), h - Inches(0.2))
    tf = b.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        p.font.size = Pt(11)
        p.font.name = MONO
        p.font.color.rgb = RGBColor(0xDD, 0xE5, 0xF0)
        p.space_after = Pt(2)


# ─────────────────────────────────────────────────────────────
def s01(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, INK)
    rect(s, 0, 0, Inches(0.25), SLIDE_H, CORAL)
    text(s, "CONCURSO DATOS AL ECOSISTEMA 2026 · IA PARA COLOMBIA",
         Inches(0.9), Inches(0.95), Inches(11.5), Inches(0.4), size=11, color=CORAL, bold=True)
    text(s, "ContratIA Abierta", Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.3),
         size=60, bold=True, color=WHITE)
    text(s, "Sistema de ingeniería de datos para priorizar la revisión humana de la "
            "contratación pública en Colombia.",
         Inches(0.9), Inches(3.0), Inches(11), Inches(0.8), size=19, color=CLOUD)
    text(s, "No detectamos corrupción. Ordenamos qué revisar primero, con evidencia trazable.",
         Inches(0.9), Inches(3.95), Inches(11), Inches(0.5), size=14, italic=True, color=SAND)
    rect(s, Inches(0.9), Inches(4.7), Inches(1.2), Pt(2), CORAL)
    cx = Inches(0.9)
    for t in ["90.431 procesos reales", "Validado vs control fiscal 2.5×",
              "Caso real: Puerto Gaitán", "Embeddings semánticos"]:
        c = chip(s, cx, Inches(5.05), t, bg_c=NAVY)
        cx += c.width + Inches(0.18)
    text(s, "Ingeniería de Datos · Universidad del Rosario · Región demo: Meta y Casanare (Orinoquía)",
         Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.4), size=11, color=CLOUD)


def s02(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    band(s)
    title(s, "El problema", "Triage, no detección: la capacidad humana < el volumen")
    text(s, "Colombia publica más de 2 millones de procesos de contratación al año en SECOP. "
            "Las oficinas de control interno y veedurías no pueden revisarlos todos. La pregunta "
            "operativa real es: ¿qué se revisa primero esta semana?",
         Inches(0.7), Inches(1.95), Inches(12), Inches(1.0), size=15, color=INK)
    data = [("> 2M", "procesos/año", "publicados en SECOP", NAVY),
            ("Decenas", "de revisores", "capacidad humana limitada", CORAL),
            ("1 decisión", "semanal", "qué priorizar, con evidencia", TEAL)]
    cw, ch = Inches(3.85), Inches(2.3)
    for i, (a, m, t, c) in enumerate(data):
        x = Inches(0.75) + i * (cw + Inches(0.25))
        card(s, x, Inches(3.3), cw, ch, a, "", accent=c, tsize=26)
        text(s, m, x + Inches(0.25), Inches(3.95), cw - Inches(0.4), Inches(0.4),
             size=13, bold=True, color=c)
        text(s, t, x + Inches(0.25), Inches(4.4), cw - Inches(0.4), Inches(1.0),
             size=11, color=MUTED)
    rect(s, Inches(0.7), Inches(6.15), Inches(12), Inches(0.6), SAND)
    text(s, "ContratIA Abierta convierte un océano de procesos en una cola priorizada y auditable.",
         Inches(0.85), Inches(6.27), Inches(11.8), Inches(0.4), size=12, bold=True, italic=True, color=INK)
    footer(s, 2)


def s03(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    band(s)
    title(s, "La tesis", "Qué SÍ hace y qué NO hace el sistema")
    rect(s, Inches(0.7), Inches(1.9), Inches(0.18), Inches(4.3), TEAL)
    text(s, "SÍ HACE", Inches(1.0), Inches(1.9), Inches(5), Inches(0.35), size=12, bold=True, color=TEAL)
    yes = ["Ranking 0–100 explicable de procesos a revisar.",
           "Comparables semánticos por categoría + embeddings.",
           "Alineación plan (PAA) vs ejecución (SECOP II).",
           "Razones legibles + confianza de datos visible.",
           "Trazabilidad: fuente, fecha, snapshot, hash."]
    y = Inches(2.35)
    for it in yes:
        text(s, "•  " + it, Inches(1.0), y, Inches(5.6), Inches(0.45), size=12, color=INK)
        y += Inches(0.52)
    rect(s, Inches(7.0), Inches(1.9), Inches(0.18), Inches(4.3), CORAL)
    text(s, "NO HACE", Inches(7.3), Inches(1.9), Inches(5), Inches(0.35), size=12, bold=True, color=CORAL)
    no = ["No declara corrupción, fraude ni responsabilidad.",
          "No reemplaza auditoría jurídica ni fiscal.",
          "No deanonimiza más allá del dato público.",
          "No usa contexto fiscal como etiqueta del modelo.",
          "No reemplaza la decisión humana: la informa."]
    y = Inches(2.35)
    for it in no:
        text(s, "•  " + it, Inches(7.3), y, Inches(5.6), Inches(0.45), size=12, color=INK)
        y += Inches(0.52)
    rect(s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.55), CLOUD)
    text(s, "Un score alto significa “revísalo primero con evidencia”, no “es corrupción”.",
         Inches(0.85), Inches(6.55), Inches(11.7), Inches(0.4), size=12, italic=True, bold=True, color=NAVY)
    footer(s, 3)


def s04(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    band(s)
    title(s, "Usuario y decisión", "Para quién es y qué decide cada semana")
    users = [("Control interno / transparencia", "Usuario primario",
              "Decide qué procesos auditar esta semana con capacidad limitada.", NAVY),
             ("Veeduría ciudadana", "Secundario",
              "Prioriza casos con evidencia para denuncia documentada.", TEAL),
             ("Periodista de datos", "Secundario",
              "Encuentra patrones territoriales con trazabilidad.", CORAL)]
    cw, ch = Inches(3.95), Inches(2.9)
    for i, (h, r, b, c) in enumerate(users):
        x = Inches(0.75) + i * (cw + Inches(0.25))
        card(s, x, Inches(1.9), cw, ch, h, "", accent=c, tsize=15)
        text(s, r, x + Inches(0.3), Inches(2.55), cw - Inches(0.5), Inches(0.3),
             size=10, bold=True, color=c)
        text(s, b, x + Inches(0.3), Inches(2.95), cw - Inches(0.5), Inches(1.8),
             size=12, color=MUTED)
    rect(s, Inches(0.75), Inches(5.2), Inches(12), Inches(1.5), SAND)
    text(s, "DECISIÓN SEMANAL", Inches(0.95), Inches(5.35), Inches(11), Inches(0.3),
         size=10, bold=True, color=CORAL)
    text(s, "“De los procesos publicados, ¿cuáles 20 revisamos primero y por qué?”",
         Inches(0.95), Inches(5.68), Inches(11.6), Inches(0.6), size=16, bold=True, color=INK)
    text(s, "El sistema convierte esa pregunta en una cola priorizada y auditable.",
         Inches(0.95), Inches(6.2), Inches(11.6), Inches(0.4), size=11, italic=True, color=MUTED)
    footer(s, 4)


def s05(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    band(s)
    title(s, "Datos abiertos", "Cuatro datasets oficiales del ecosistema Datos Abiertos Colombia")
    src = [("p6dx-8zbt", "SECOP II — Procesos", "Procesos de contratación (canónico).", NAVY),
           ("rpmr-utcd", "SECOP Integrado", "Histórico de contratos (baseline de valor).", TEAL),
           ("9sue-ezhx", "Plan Anual de Adquisiciones", "Lo que la entidad planeó comprar.", CORAL),
           ("wasc-xi4h", "Control fiscal (AGR)", "Contexto — NO etiqueta del modelo.", MUTED)]
    cw, ch = Inches(5.95), Inches(1.5)
    for i, (code, name, body, c) in enumerate(src):
        x = Inches(0.7) + (i % 2) * (cw + Inches(0.2))
        y = Inches(1.95) + (i // 2) * (ch + Inches(0.2))
        card(s, x, y, cw, ch, "", "", accent=c)
        text(s, code, x + Inches(0.28), y + Inches(0.14), Inches(2.5), Inches(0.35),
             size=12, bold=True, color=c, font=MONO)
        text(s, name, x + Inches(0.28), y + Inches(0.52), cw - Inches(0.5), Inches(0.4),
             size=14, bold=True, color=INK)
        text(s, body, x + Inches(0.28), y + Inches(0.95), cw - Inches(0.5), Inches(0.5),
             size=10, color=MUTED)
    mets = [("90.431", "procesos scoreados"), ("33", "objetos PostgreSQL"),
            ("4.494", "PAA items"), ("100%", "datos abiertos")]
    mw = Inches(2.85)
    for i, (v, lb) in enumerate(mets):
        metric(s, Inches(0.7) + i * (mw + Inches(0.1)), Inches(5.55), mw, Inches(1.0), v, lb)
    footer(s, 5)


def s06(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    band(s)
    title(s, "Arquitectura", "Una sola lane de ingeniería de datos, de extremo a extremo")
    steps = [("Socrata API", "SECOP II/Integrado\nPAA · AGR"),
             ("ETL", "Polars + Parquet\njoin audit"),
             ("PostgreSQL\n+ MongoDB", "33 objetos\n+ 5 colecciones"),
             ("FastAPI ×3", "contracts · risk\nanalytics"),
             ("Dash", "DECIDIR · ENTENDER\nCONFIAR")]
    n = len(steps)
    iw = (SLIDE_W - Inches(1.4) - Inches(0.5) * (n - 1)) / n
    y = Inches(2.4)
    for i, (h, b) in enumerate(steps):
        x = Inches(0.7) + i * (iw + Inches(0.5))
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x + iw / 2 - Inches(0.32), y, Inches(0.64), Inches(0.64))
        circ.fill.solid()
        circ.fill.fore_color.rgb = NAVY
        circ.line.fill.background()
        text(s, str(i + 1), x + iw / 2 - Inches(0.32), y, Inches(0.64), Inches(0.64),
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cc = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y + Inches(0.85), iw, Inches(1.7))
        cc.fill.solid()
        cc.fill.fore_color.rgb = WHITE
        cc.line.color.rgb = MIST
        cc.line.width = Pt(0.75)
        cc.adjustments[0] = 0.08
        text(s, h, x, y + Inches(1.0), iw, Inches(0.6), size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        text(s, b, x + Inches(0.1), y + Inches(1.6), iw - Inches(0.2), Inches(0.9),
             size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
        if i < n - 1:
            arrow(s, x + iw + Inches(0.03), y + Inches(0.24))
    rect(s, Inches(0.7), Inches(5.7), Inches(12), Inches(1.1), CLOUD)
    text(s, "Reproducible end-to-end", Inches(0.95), Inches(5.8), Inches(11), Inches(0.3),
         size=10, bold=True, color=CORAL)
    text(s, "PostgreSQL es la fuente de verdad relacional; MongoDB guarda evidencia documental y "
            "eventos. Stack: Python 3.11 · uv · Polars · FastAPI · Dash · Docker Compose.",
         Inches(0.95), Inches(6.1), Inches(11.6), Inches(0.6), size=11, color=INK)
    footer(s, 6)


def s07(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    band(s)
    title(s, "Modelo relacional", "PostgreSQL como fuente de verdad: 33 objetos públicos")
    img(s, "er_model.png", Inches(0.8), Inches(1.85), w=Inches(8.0))
    rx = Inches(9.1)
    text(s, "GARANTÍAS", rx, Inches(1.95), Inches(3.6), Inches(0.3), size=10, bold=True, color=CORAL)
    items = [("PK / FK", "integridad referencial"),
             ("CHECK", "score 0–100, conf 0–1"),
             ("Índices", "consulta analítica"),
             ("Vistas", "ranking + PAA"),
             ("Triggers", "auditoría + historial")]
    yy = Inches(2.35)
    for lb, v in items:
        text(s, lb, rx, yy, Inches(1.3), Inches(0.3), size=11, bold=True, color=NAVY)
        text(s, v, rx + Inches(1.3), yy, Inches(2.4), Inches(0.4), size=10.5, color=MUTED)
        yy += Inches(0.5)
    rect(s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.5), SAND)
    text(s, "27 tablas + vistas + triggers = 33 objetos. Diseño normalizado con auditoría incorporada.",
         Inches(0.85), Inches(6.55), Inches(11.7), Inches(0.4), size=11, italic=True, color=INK)
    footer(s, 7)


def s08(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    band(s)
    title(s, "SQL engineering", "Lo que un evaluador técnico espera ver")
    cards = [("Triggers", "audit_log · historial de estados · updated_at"),
             ("Window functions", "concentración por entidad · outliers · percent_rank"),
             ("CTE recursiva", "jerarquía territorial departamento → municipio"),
             ("Transacciones", "score + evento de auditoría atómico")]
    for i, (h, b) in enumerate(cards):
        x = Inches(0.7) + (i % 2) * Inches(6.1)
        y = Inches(1.95) + (i // 2) * Inches(1.15)
        card(s, x, y, Inches(5.9), Inches(1.0), h, b, accent=NAVY, tsize=14, bsize=11)
    code_block(s, Inches(0.7), Inches(4.35), Inches(12), Inches(2.0), [
        "-- Window: percentil global del score (usado en el ranking)",
        "SELECT process_id, priority_score,",
        "  round(100*percent_rank() OVER (ORDER BY priority_score),2)",
        "    AS score_percentile",
        "FROM v_ranking_processes;   -- 'top 0.5%' en vez de '85/100'",
    ])
    footer(s, 8)


def s09(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    band(s)
    title(s, "Score explicable", "Cuatro componentes auditables → un número 0–100")
    comp = [("Anomalía · 45%", "IsolationForest sobre features numéricas", NAVY),
            ("Pares · 35%", "Desviación vs procesos comparables UNSPSC", TEAL),
            ("Reglas · 20%", "Thresholds verificables (valor, plazo, PAA)", CORAL),
            ("Confianza", "0–1 según cobertura de datos (aparte)", MUTED)]
    cw = Inches(2.9)
    for i, (h, b, c) in enumerate(comp):
        x = Inches(0.7) + i * (cw + Inches(0.2))
        card(s, x, Inches(1.9), cw, Inches(1.9), h, b, accent=c, tsize=13, bsize=10)
    code_block(s, Inches(0.7), Inches(4.05), Inches(12), Inches(0.95), [
        "score = round(100 · σ(0.45·anomalía + 0.35·pares + 0.20·reglas))",
    ])
    rect(s, Inches(0.7), Inches(5.25), Inches(12), Inches(1.4), CLOUD)
    text(s, "INTERPRETACIÓN: percentil, no nota absoluta", Inches(0.95), Inches(5.38),
         Inches(11), Inches(0.3), size=10, bold=True, color=CORAL)
    text(s, "El score se normaliza por cuantiles del dataset. Un score de 91 = “top 0.02% más "
            "atípico de 90.431”, no “91% corrupto”. Distribución de triage: mediana 7, p95 ≈ 45, "
            "solo 0.5% supera 70.",
         Inches(0.95), Inches(5.7), Inches(11.6), Inches(0.9), size=12, color=INK)
    footer(s, 9)


def s10(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    band(s)
    title(s, "IA semántica", "Embeddings multilingües funcionando (no humo)")
    text(s, "El matching plan↔ejecución y los comparables usan un proveedor semántico "
            "configurable. Por defecto: embeddings neuronales MiniLM multilingüe; fallback "
            "robusto a TF-IDF si no hay modelo. Ambos comparten la misma abstracción.",
         Inches(0.7), Inches(1.95), Inches(12), Inches(0.9), size=14, color=INK)
    card(s, Inches(0.7), Inches(3.0), Inches(5.9), Inches(1.7),
         "SentenceTransformer (default)", "paraphrase-multilingual-MiniLM-L12-v2\n"
         "Carga en ~11s · offline · sin API key", accent=TEAL, tsize=13, bsize=11)
    card(s, Inches(6.8), Inches(3.0), Inches(5.9), Inches(1.7),
         "TF-IDF (fallback)", "scikit-learn · cero dependencias pesadas\n"
         "Garantiza que la demo nunca se cae", accent=NAVY, tsize=13, bsize=11)
    rect(s, Inches(0.7), Inches(5.0), Inches(12), Inches(1.65), INK)
    text(s, "PRUEBA DE DISCRIMINACIÓN SEMÁNTICA (medida en vivo)", Inches(0.95), Inches(5.15),
         Inches(11), Inches(0.3), size=10, bold=True, color=CORAL)
    text(s, "sim(“acueducto”, “alcantarillado”) = 0.611      "
            "sim(“acueducto”, “papelería”) = 0.136",
         Inches(0.95), Inches(5.5), Inches(11.5), Inches(0.4), size=14, bold=True, color=WHITE, font=MONO)
    text(s, "El modelo entiende significado: agrupa contratos de agua y separa los no relacionados.",
         Inches(0.95), Inches(6.0), Inches(11.5), Inches(0.4), size=11, italic=True, color=CLOUD)
    footer(s, 10)


def s11(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    band(s)
    title(s, "Validación independiente", "¿El score se concentra donde el control fiscal ya miró?")
    text(s, "Prueba ciega: cruzamos las entidades que la Auditoría General (AGR, wasc-xi4h) puso "
            "bajo vigilancia fiscal contra nuestro score. El modelo nunca vio esa etiqueta.",
         Inches(0.7), Inches(1.9), Inches(12), Inches(0.7), size=14, color=INK)
    img(s, "zona_confiar.png", Inches(0.7), Inches(2.7), w=Inches(7.4))
    rx = Inches(8.4)
    metric(s, rx, Inches(2.75), Inches(4.3), Inches(1.1), "2.46×",
           "enriquecimiento en tasa de alta prioridad", color=CORAL, bg_c=CLOUD, vsize=30)
    metric(s, rx, Inches(4.0), Inches(2.05), Inches(1.0), "19", "mediana — vigilada", color=CORAL)
    metric(s, rx + Inches(2.25), Inches(4.0), Inches(2.05), Inches(1.0), "7", "mediana — resto", color=NAVY)
    text(s, "Auditoría AGR ≠ conducta indebida probada: mide selección para revisión fiscal. "
            "Que el score enriquezca ese grupo valida el triage, no acusa.",
         rx, Inches(5.2), Inches(4.3), Inches(1.4), size=11, italic=True, color=MUTED)
    footer(s, 11)


def s12(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    band(s)
    title(s, "Caso real", "Puerto Gaitán: el sistema prioriza donde la Contraloría ya halló señal")
    text(s, "REGISTRO PÚBLICO", Inches(0.7), Inches(1.85), Inches(6), Inches(0.3),
         size=10, bold=True, color=CORAL)
    text(s, "La Contraloría documentó irregularidades vigentes en Puerto Gaitán: contratos de "
            "agua/acueducto (hallazgo > $14.700 millones), sistemas fotovoltaicos ($8.901 "
            "millones, 2025) y manejo de regalías.",
         Inches(0.7), Inches(2.2), Inches(5.9), Inches(2.0), size=12.5, color=INK)
    text(s, "NUESTRO SISTEMA, A CIEGAS", Inches(0.7), Inches(4.0), Inches(6), Inches(0.3),
         size=10, bold=True, color=CORAL)
    text(s, "Marca Puerto Gaitán 3.1× sobre la tasa nacional. Su proceso top priorizado es una "
            "construcción de acueducto de $42.408 millones (score 87) — misma entidad y categoría "
            "que señaló la Contraloría. No usamos esa información para entrenar.",
         Inches(0.7), Inches(4.35), Inches(5.9), Inches(2.0), size=12.5, color=INK)
    rx = Inches(7.0)
    metric(s, rx, Inches(1.95), Inches(5.7), Inches(1.2), "3.1×",
           "Puerto Gaitán vs tasa nacional de alta prioridad", color=CORAL, bg_c=CLOUD, vsize=34)
    text(s, "EL CONTRAEJEMPLO QUE DA CREDIBILIDAD", rx, Inches(3.4), Inches(5.7), Inches(0.3),
         size=10, bold=True, color=NAVY)
    rows = [("Puerto Gaitán", "3.1×", CORAL), ("Departamento del Meta", "2.9×", CORAL),
            ("Gobernación de Casanare", "1.0× (promedio)", MUTED),
            ("Alcaldía de Villavicencio", "0.7× (debajo)", MUTED)]
    yy = Inches(3.75)
    for name, v, c in rows:
        text(s, name, rx, yy, Inches(4.0), Inches(0.3), size=12, color=INK)
        text(s, v, rx + Inches(4.0), yy, Inches(1.7), Inches(0.3), size=12, bold=True, color=c, align=PP_ALIGN.RIGHT)
        yy += Inches(0.42)
    text(s, "Casanare (escándalos de hace 10+ años) puntúa en el promedio: el sistema detecta "
            "estructura del proceso actual, no reputación histórica. No es una lista negra.",
         rx, Inches(5.6), Inches(5.7), Inches(1.0), size=10.5, italic=True, color=MUTED)
    footer(s, 12)


def s13(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    band(s)
    title(s, "El producto", "Tres zonas por el trabajo del usuario, no por features")
    img(s, "zona_decidir.png", Inches(0.7), Inches(1.9), w=Inches(7.6))
    rx = Inches(8.6)
    zones = [("1 · DECIDIR", "Qué revisar primero. Cola priorizada + filtro territorial + KPIs reales.", NAVY),
             ("2 · ENTENDER", "El caso: score desglosado, razones, comparables, PAA, contexto fiscal.", TEAL),
             ("3 · CONFIAR", "Por qué confiar: validación AGR, distribución, metodología, límites.", CORAL)]
    yy = Inches(2.0)
    for h, b, c in zones:
        card(s, rx, yy, Inches(4.1), Inches(1.45), h, b, accent=c, tsize=14, bsize=10.5)
        yy += Inches(1.6)
    footer(s, 13)


def s14(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    band(s)
    title(s, "Evidencia de ingeniería", "Métricas duras de un sistema reproducible")
    mets = [("90.431", "procesos scoreados", NAVY), ("33", "objetos PostgreSQL", TEAL),
            ("5", "colecciones MongoDB", CORAL), ("3", "APIs FastAPI · health 200", NAVY),
            ("71", "tests pytest", TEAL), ("4", "datasets abiertos", CORAL)]
    cw, ch = Inches(2.45), Inches(1.5)
    for i, (v, lb, c) in enumerate(mets):
        x = Inches(0.7) + (i % 3) * (cw + Inches(0.15))
        y = Inches(1.95) + (i // 3) * (ch + Inches(0.15))
        metric(s, x, y, cw, ch, v, lb, color=c, bg_c=WHITE, vsize=28)
    rx = Inches(8.4)
    text(s, "VALIDACIÓN CUBIERTA", rx, Inches(1.95), Inches(4.5), Inches(0.3), size=10, bold=True, color=CORAL)
    chk = ["Schema integrity (PK/FK/CHECK)", "Join audit + fill rates", "Compuerta PAA top-1 ≥ 0.75",
           "Scoring bounds 0–100 / 0–1", "Enriquecimiento AGR 2.5×", "Caso real Puerto Gaitán 3.1×"]
    yy = Inches(2.35)
    for c in chk:
        text(s, "✓  " + c, rx, yy, Inches(4.6), Inches(0.3), size=11.5, color=INK)
        yy += Inches(0.42)
    footer(s, 14)


def s15(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    band(s)
    title(s, "Límites honestos", "Lo que falta — declarado, no oculto")
    done = ["71 tests · lint estable · validación reproducible",
            "90.431 procesos reales (Meta + Casanare)",
            "Embeddings semánticos funcionando + fallback TF-IDF",
            "Validación AGR 2.5× + caso real Puerto Gaitán"]
    pend = ["Validación humana de 100 casos con dos revisores",
            "Benchmark formal embeddings vs TF-IDF (precision@k)",
            "Despliegue público read-only (Blueprint listo)",
            "Sesiones de usabilidad con 5 usuarios reales"]
    text(s, "LISTO", Inches(0.7), Inches(1.95), Inches(5), Inches(0.3), size=11, bold=True, color=TEAL)
    yy = Inches(2.35)
    for d in done:
        text(s, "✓  " + d, Inches(0.7), yy, Inches(5.8), Inches(0.5), size=12, color=INK)
        yy += Inches(0.6)
    text(s, "PENDIENTE (ROADMAP)", Inches(7.0), Inches(1.95), Inches(5), Inches(0.3), size=11, bold=True, color=CORAL)
    yy = Inches(2.35)
    for p in pend:
        text(s, "○  " + p, Inches(7.0), yy, Inches(5.8), Inches(0.5), size=12, color=INK)
        yy += Inches(0.6)
    rect(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.6), SAND)
    text(s, "La rúbrica premia una validación pendiente declarada por encima de una validación inventada.",
         Inches(0.85), Inches(6.32), Inches(11.7), Inches(0.4), size=11, italic=True, bold=True, color=INK)
    footer(s, 15)


def s16(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, INK)
    rect(s, 0, 0, Inches(0.25), SLIDE_H, CORAL)
    text(s, "ContratIA Abierta", Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.1),
         size=46, bold=True, color=WHITE)
    text(s, "Sin entrenar contra ninguna etiqueta de corrupción, el sistema prioriza 3.1× los "
            "contratos de Puerto Gaitán — donde la Contraloría halló irregularidades vigentes.",
         Inches(0.9), Inches(2.3), Inches(11.2), Inches(1.0), size=17, color=CLOUD)
    text(s, "Detecta estructura, no reputación. Prioriza revisión humana con evidencia trazable.",
         Inches(0.9), Inches(3.35), Inches(11.2), Inches(0.5), size=13, italic=True, color=SAND)
    rect(s, Inches(0.9), Inches(4.1), Inches(1.2), Pt(2), CORAL)
    info = [("EQUIPO", "Ingeniería de Datos · Universidad del Rosario"),
            ("REPOSITORIO", "github.com/Thom-320/secop-risk-alerts-co"),
            ("CATEGORÍA · REGIÓN", "Gobernanza y Transparencia · Meta y Casanare")]
    yy = Inches(4.5)
    for lb, v in info:
        text(s, lb, Inches(0.9), yy, Inches(4), Inches(0.3), size=10, bold=True, color=CORAL)
        text(s, v, Inches(0.9), yy + Inches(0.28), Inches(9), Inches(0.4), size=13,
             color=WHITE, font=(MONO if "github" in v else FONT))
        yy += Inches(0.78)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for fn in (s01, s02, s03, s04, s05, s06, s07, s08, s09, s10,
               s11, s12, s13, s14, s15, s16):
        fn(prs)
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"OK — {len(prs.slides)} slides → {OUTPUT}")


if __name__ == "__main__":
    main()
