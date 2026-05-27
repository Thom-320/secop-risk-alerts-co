#!/usr/bin/env python3
"""ContratIA Abierta — Deck v3 (concurso-grade, ES)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ── Canvas ─────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Palette (deliberate, professional) ─────────────────────────
INK       = RGBColor(0x0B, 0x1E, 0x33)  # titles, primary text
NAVY      = RGBColor(0x10, 0x3A, 0x5C)  # primary accent
CORAL     = RGBColor(0xE3, 0x5A, 0x4B)  # call-out / highlights
TEAL      = RGBColor(0x1F, 0x82, 0x7C)  # positive / metrics
SAND      = RGBColor(0xF5, 0xEA, 0xD2)  # warm soft bg
PAPER     = RGBColor(0xFA, 0xFA, 0xF7)  # canvas
CLOUD     = RGBColor(0xEE, 0xF2, 0xF8)  # soft chip bg
MIST      = RGBColor(0xDB, 0xE2, 0xEC)  # rule lines
MUTED     = RGBColor(0x55, 0x60, 0x76)  # secondary text
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DEEP      = RGBColor(0x05, 0x12, 0x22)  # near-black for hero

FONT = "Calibri"
ASSETS = Path(__file__).resolve().parent.parent / "assets"
OUTPUT = Path(__file__).resolve().parent.parent / "contratia_abierta_deck_v3.pptx"

FOOTER = "ContratIA Abierta  ·  Datos al Ecosistema 2026  ·  Prioriza revisión humana, no acusa."


# ── Primitives ─────────────────────────────────────────────────
def set_bg(slide, color=PAPER):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if not line:
        shp.line.fill.background()
    return shp


def add_top_band(slide, color=NAVY, accent=CORAL):
    add_rect(slide, 0, 0, SLIDE_W, Pt(6), color)
    add_rect(slide, 0, Pt(6), Inches(1.6), Pt(3), accent)


def add_text(slide, text, left, top, width, height,
             size=14, color=INK, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font
    return box


def add_multi_text(slide, lines, left, top, width, height, font=FONT, align=PP_ALIGN.LEFT):
    """lines: list of dicts with text, size, color, bold, italic, space_after."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line.get("text", "")
        p.alignment = line.get("align", align)
        p.font.size = Pt(line.get("size", 14))
        p.font.bold = line.get("bold", False)
        p.font.italic = line.get("italic", False)
        p.font.color.rgb = line.get("color", INK)
        p.font.name = font
        if "space_after" in line:
            p.space_after = Pt(line["space_after"])
        if "space_before" in line:
            p.space_before = Pt(line["space_before"])
    return box


def add_title(slide, eyebrow, title, top=Inches(0.55)):
    add_text(slide, eyebrow.upper(), Inches(0.7), top,
             Inches(10), Inches(0.3), size=10, color=CORAL, bold=True)
    add_text(slide, title, Inches(0.7), top + Inches(0.32),
             Inches(12), Inches(0.7), size=26, bold=True, color=INK)
    add_rect(slide, Inches(0.7), top + Inches(1.05),
             Inches(0.6), Pt(2), CORAL)


def add_footer(slide, page=None, total=None):
    add_text(slide, FOOTER, Inches(0.7), SLIDE_H - Inches(0.4),
             Inches(10), Inches(0.3), size=8, color=MUTED, italic=True)
    if page is not None:
        pg = f"{page} / {total}" if total else str(page)
        add_text(slide, pg, SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.4),
                 Inches(0.6), Inches(0.3), size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_card(slide, left, top, width, height, title, body,
             title_color=INK, body_color=MUTED, bg=WHITE, border=MIST,
             title_size=13, body_size=10, accent_top=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = bg
    shp.line.color.rgb = border
    shp.line.width = Pt(0.75)
    shp.adjustments[0] = 0.04

    if accent_top:
        add_rect(slide, left, top, width, Pt(3), accent_top)

    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12); tf.margin_right = Pt(12)
    tf.margin_top = Pt(10); tf.margin_bottom = Pt(10)

    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.size = Pt(title_size)
    p1.font.bold = True
    p1.font.color.rgb = title_color
    p1.font.name = FONT
    p1.space_after = Pt(4)

    if body:
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(body_size)
        p2.font.color.rgb = body_color
        p2.font.name = FONT


def add_metric(slide, left, top, width, height, value, label,
               value_color=NAVY, bg=CLOUD, value_size=24):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = bg
    shp.line.fill.background()
    shp.adjustments[0] = 0.08

    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top = Pt(8); tf.margin_bottom = Pt(8)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = tf.paragraphs[0]
    p1.text = value
    p1.font.size = Pt(value_size)
    p1.font.bold = True
    p1.font.color.rgb = value_color
    p1.font.name = FONT
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(9)
    p2.font.color.rgb = MUTED
    p2.font.name = FONT
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)


def add_chip(slide, left, top, text, bg=CLOUD, fg=NAVY, width=None):
    w = width or Inches(max(1.2, len(text) * 0.09 + 0.4))
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, w, Inches(0.34))
    shp.fill.solid()
    shp.fill.fore_color.rgb = bg
    shp.line.fill.background()
    shp.adjustments[0] = 0.5

    tf = shp.text_frame
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top = Pt(3); tf.margin_bottom = Pt(3)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = fg
    p.font.name = FONT
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return shp


def add_arrow(slide, left, top, width=Inches(0.5), color=NAVY):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.22))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()


def add_image_safe(slide, name, left, top, width=None, height=None):
    path = ASSETS / name
    if not path.exists():
        add_card(slide, left, top, width or Inches(4), height or Inches(2.4),
                 f"[{name}]", "Imagen no encontrada", bg=CLOUD)
        return None
    kwargs = {}
    if width: kwargs["width"] = width
    if height: kwargs["height"] = height
    return slide.shapes.add_picture(str(path), left, top, **kwargs)


def add_rule(slide, left, top, width, color=MIST, thickness=Pt(1)):
    return add_rect(slide, left, top, width, thickness, color)


# ── Slides ─────────────────────────────────────────────────────
TOTAL = 13


def s01_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, INK)

    # accent corners
    add_rect(slide, 0, 0, Inches(0.25), SLIDE_H, CORAL)
    add_rect(slide, SLIDE_W - Inches(2.6), 0, Inches(2.6), Inches(0.18), CORAL)

    # eyebrow
    add_text(slide, "CONCURSO DATOS AL ECOSISTEMA 2026  ·  IA PARA COLOMBIA",
             Inches(0.9), Inches(1.0), Inches(11), Inches(0.4),
             size=11, color=CORAL, bold=True)

    # brand
    add_text(slide, "ContratIA Abierta",
             Inches(0.9), Inches(1.55), Inches(11), Inches(1.4),
             size=64, bold=True, color=WHITE)

    # subtitle
    add_text(slide, "IA explicable para priorizar la revisión humana de la contratación pública en Colombia.",
             Inches(0.9), Inches(3.2), Inches(11), Inches(0.7),
             size=20, color=CLOUD)

    # tagline
    add_text(slide, "No emitimos juicios de corrupción.  Ordenamos qué revisar primero, con evidencia trazable.",
             Inches(0.9), Inches(4.1), Inches(11), Inches(0.5),
             size=14, italic=True, color=SAND)

    add_rect(slide, Inches(0.9), Inches(4.85), Inches(1.2), Pt(2), CORAL)

    # tags
    cx = Inches(0.9)
    for tag in ["Datos abiertos SECOP", "Score 0–100 explicable", "Revisión humana trazable", "Orinoquía → Colombia"]:
        chip = add_chip(slide, cx, Inches(5.2), tag, bg=NAVY, fg=WHITE)
        cx += chip.width + Inches(0.18)

    # bottom strip
    add_text(slide, "Categoría: Gobernanza y Transparencia   ·   Región demo: Meta y Casanare   ·   Universidad del Rosario",
             Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.4),
             size=11, color=CLOUD)


def s02_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide)
    add_title(slide, "El problema", "No faltan datos. Falta capacidad humana para revisarlos.")

    add_text(slide,
             "Cada año Colombia publica cientos de miles de procesos en SECOP. Oficinas de "
             "control interno, veedurías y periodistas no tienen tiempo para revisarlos todos. "
             "La pregunta operativa real es: ¿qué se revisa primero esta semana?",
             Inches(0.7), Inches(1.9), Inches(12), Inches(1.2),
             size=15, color=INK)

    # 3 stat cards
    stats = [
        ("Miles", "de procesos publicados", "por entidad y año en SECOP."),
        ("Decenas", "de revisores humanos", "con tiempo y capacidad limitada."),
        ("Una decisión", "semanal real", "qué priorizar, con qué evidencia."),
    ]
    cw = Inches(3.85); ch = Inches(2.4)
    start = Inches(0.75); gap = Inches(0.25)
    for i, (head, mid, tail) in enumerate(stats):
        x = start + i * (cw + gap)
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.4), cw, ch)
        shp.fill.solid(); shp.fill.fore_color.rgb = WHITE
        shp.line.color.rgb = MIST; shp.line.width = Pt(0.75); shp.adjustments[0] = 0.04
        add_rect(slide, x, Inches(3.4), Pt(4), ch, [NAVY, CORAL, TEAL][i])
        add_text(slide, head, x + Inches(0.25), Inches(3.55),
                 cw - Inches(0.4), Inches(0.6), size=26, bold=True, color=INK)
        add_text(slide, mid, x + Inches(0.25), Inches(4.2),
                 cw - Inches(0.4), Inches(0.4), size=13, bold=True, color=NAVY)
        add_text(slide, tail, x + Inches(0.25), Inches(4.65),
                 cw - Inches(0.4), Inches(1.2), size=11, color=MUTED)

    # insight strip
    add_rect(slide, Inches(0.7), Inches(6.2), Inches(12), Inches(0.6), SAND)
    add_text(slide,
             "Triage de contratación: convertir un océano de procesos en una cola priorizada y auditable.",
             Inches(0.85), Inches(6.32), Inches(11.8), Inches(0.4),
             size=12, bold=True, color=INK, italic=True)

    add_footer(slide, 2, TOTAL)


def s03_thesis(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide)
    add_title(slide, "La tesis", "Qué hace ContratIA Abierta y qué NO hace")

    # 2 columns
    # YES column
    add_rect(slide, Inches(0.7), Inches(1.85), Inches(0.18), Inches(4.5), TEAL)
    add_text(slide, "SÍ HACE", Inches(1.0), Inches(1.85),
             Inches(5), Inches(0.4), size=11, bold=True, color=TEAL)
    yes_items = [
        "Ranking 0–100 explicable de procesos a revisar.",
        "Comparables: ¿con qué se parece este proceso?",
        "Alineación plan (PAA) vs ejecución (SECOP II).",
        "Razones legibles para humano + confianza visible.",
        "Trazabilidad: fuente, fecha, snapshot, hash.",
    ]
    y = Inches(2.3)
    for it in yes_items:
        add_text(slide, "•  " + it, Inches(1.0), y, Inches(5.6), Inches(0.45),
                 size=12, color=INK)
        y += Inches(0.55)

    # NO column
    add_rect(slide, Inches(7.0), Inches(1.85), Inches(0.18), Inches(4.5), CORAL)
    add_text(slide, "NO HACE", Inches(7.3), Inches(1.85),
             Inches(5), Inches(0.4), size=11, bold=True, color=CORAL)
    no_items = [
        "No declara corrupción, fraude ni responsabilidad.",
        "No reemplaza auditoría jurídica ni fiscal.",
        "No deanonimiza más allá del dato público.",
        "No usa contexto fiscal como etiqueta del modelo.",
        "No reemplaza la decisión humana: la informa.",
    ]
    y = Inches(2.3)
    for it in no_items:
        add_text(slide, "•  " + it, Inches(7.3), y, Inches(5.6), Inches(0.45),
                 size=12, color=INK)
        y += Inches(0.55)

    add_rect(slide, Inches(0.7), Inches(6.55), Inches(12), Inches(0.5), CLOUD)
    add_text(slide,
             "Un score alto significa “revísalo primero con evidencia”, no “es corrupción”.",
             Inches(0.85), Inches(6.62), Inches(11.7), Inches(0.4),
             size=12, italic=True, bold=True, color=NAVY)

    add_footer(slide, 3, TOTAL)


def s04_user(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide)
    add_title(slide, "Usuario y decisión", "¿Para quién es y qué decide cada semana?")

    users = [
        ("Control interno / transparencia",
         "Usuario primario.",
         "Decide qué procesos auditar esta semana con capacidad limitada.",
         NAVY),
        ("Veeduría ciudadana",
         "Usuario secundario.",
         "Prioriza casos con evidencia para denuncia ciudadana documentada.",
         TEAL),
        ("Periodista de datos",
         "Usuario secundario.",
         "Encuentra patrones territoriales con trazabilidad para publicar.",
         CORAL),
    ]
    cw = Inches(3.95); ch = Inches(3.2)
    start = Inches(0.75); gap = Inches(0.25)
    for i, (head, role, body, color) in enumerate(users):
        x = start + i * (cw + gap)
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.85), cw, ch)
        shp.fill.solid(); shp.fill.fore_color.rgb = WHITE
        shp.line.color.rgb = MIST; shp.line.width = Pt(0.75); shp.adjustments[0] = 0.04
        add_rect(slide, x, Inches(1.85), cw, Pt(4), color)
        add_text(slide, head, x + Inches(0.3), Inches(2.05),
                 cw - Inches(0.5), Inches(0.6), size=15, bold=True, color=INK)
        add_text(slide, role, x + Inches(0.3), Inches(2.7),
                 cw - Inches(0.5), Inches(0.35), size=10, bold=True, color=color)
        add_text(slide, body, x + Inches(0.3), Inches(3.1),
                 cw - Inches(0.5), Inches(2), size=12, color=MUTED)

    # decision strip
    add_rect(slide, Inches(0.75), Inches(5.4), Inches(12), Inches(1.3), SAND)
    add_text(slide, "DECISIÓN SEMANAL", Inches(0.95), Inches(5.55),
             Inches(11), Inches(0.3), size=10, bold=True, color=CORAL)
    add_text(slide,
             "“De los procesos publicados esta semana, ¿cuáles 20 revisamos primero y por qué?”",
             Inches(0.95), Inches(5.85), Inches(11.6), Inches(0.6),
             size=15, bold=True, color=INK)
    add_text(slide,
             "ContratIA Abierta convierte esa pregunta en una cola priorizada y auditable.",
             Inches(0.95), Inches(6.25), Inches(11.6), Inches(0.4),
             size=11, italic=True, color=MUTED)

    add_footer(slide, 4, TOTAL)


def s05_data(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide)
    add_title(slide, "Fuentes oficiales", "Cuatro datasets abiertos del ecosistema Datos Abiertos Colombia")

    sources = [
        ("p6dx-8zbt", "SECOP II — Procesos",
         "Procesos de contratación pública (canónico).", NAVY),
        ("rpmr-utcd", "SECOP Integrado",
         "Histórico de contratos cerrados (baseline de valor).", TEAL),
        ("9sue-ezhx", "Plan Anual de Adquisiciones",
         "Lo que la entidad planeó comprar (PAA detalle).", CORAL),
        ("wasc-xi4h", "Resultados de auditoría (AGR)",
         "Contexto fiscal — visible, no etiqueta del modelo.", MUTED),
    ]
    cw = Inches(5.95); ch = Inches(1.65)
    sx = Inches(0.7); sy = Inches(1.95); gx = Inches(0.2); gy = Inches(0.2)
    for i, (code, name, body, color) in enumerate(sources):
        r = i // 2; c = i % 2
        x = sx + c * (cw + gx); y = sy + r * (ch + gy)
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
        shp.fill.solid(); shp.fill.fore_color.rgb = WHITE
        shp.line.color.rgb = MIST; shp.line.width = Pt(0.75); shp.adjustments[0] = 0.05
        add_rect(slide, x, y, Pt(5), ch, color)
        add_text(slide, code, x + Inches(0.3), y + Inches(0.15),
                 Inches(2.5), Inches(0.4), size=11, bold=True, color=color, font="Consolas")
        add_text(slide, name, x + Inches(0.3), y + Inches(0.55),
                 cw - Inches(0.5), Inches(0.45), size=14, bold=True, color=INK)
        add_text(slide, body, x + Inches(0.3), y + Inches(1.05),
                 cw - Inches(0.5), Inches(0.6), size=10, color=MUTED)

    # metrics row
    metrics = [("17.229", "procesos demo"),
               ("4.494", "PAA items"),
               ("68.916", "razones explicables"),
               ("100%", "datos abiertos")]
    mw = Inches(2.85); mh = Inches(1.0)
    mx = Inches(0.7); my = Inches(5.7)
    for i, (v, l) in enumerate(metrics):
        add_metric(slide, mx + i * (mw + Inches(0.1)), my, mw, mh, v, l,
                   value_color=NAVY, bg=CLOUD)

    add_footer(slide, 5, TOTAL)


def s06_how(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide)
    add_title(slide, "Cómo funciona", "Pipeline reproducible: del dato abierto al ranking explicable")

    steps = [
        ("Extracción",  "Socrata API\n+ snapshots Parquet"),
        ("Limpieza",    "Normalización,\nIDs estables, dedup."),
        ("Matching",    "PAA ↔ proceso por\nsimilitud + reglas"),
        ("Scoring",     "Reglas + pares\n+ anomalía + confianza"),
        ("Evidencia",   "Ranking + razones\n+ comparables + reporte"),
    ]
    n = len(steps)
    total_w = SLIDE_W - Inches(1.4)
    item_w = (total_w - Inches(0.6) * (n - 1)) / n
    y = Inches(2.4)
    h = Inches(2.2)
    for i, (head, body) in enumerate(steps):
        x = Inches(0.7) + i * (item_w + Inches(0.6))
        # circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + item_w/2 - Inches(0.35),
                                        y, Inches(0.7), Inches(0.7))
        circle.fill.solid(); circle.fill.fore_color.rgb = NAVY
        circle.line.fill.background()
        add_text(slide, str(i + 1),
                 x + item_w/2 - Inches(0.35), y, Inches(0.7), Inches(0.7),
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        # card
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x, y + Inches(0.9), item_w, h - Inches(0.9))
        shp.fill.solid(); shp.fill.fore_color.rgb = WHITE
        shp.line.color.rgb = MIST; shp.line.width = Pt(0.75); shp.adjustments[0] = 0.07
        add_text(slide, head, x, y + Inches(1.05),
                 item_w, Inches(0.4), size=13, bold=True,
                 color=NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + Inches(0.15), y + Inches(1.5),
                 item_w - Inches(0.3), Inches(1.4),
                 size=10, color=MUTED, align=PP_ALIGN.CENTER)
        # connector
        if i < n - 1:
            add_arrow(slide, x + item_w + Inches(0.05), y + Inches(0.27),
                      width=Inches(0.5), color=CORAL)

    # bottom note
    add_rect(slide, Inches(0.7), Inches(5.6), Inches(12), Inches(1.1), CLOUD)
    add_text(slide, "Reproducibilidad real",
             Inches(0.95), Inches(5.7), Inches(11), Inches(0.3),
             size=10, bold=True, color=CORAL)
    add_text(slide,
             "Un solo `make product-pipeline && make validate-product` reconstruye el ranking desde cero "
             "con datos abiertos versionados. Sin black box: cada paso deja artefacto auditable.",
             Inches(0.95), Inches(5.95), Inches(11.6), Inches(0.7),
             size=11, color=INK)

    add_footer(slide, 6, TOTAL)


def s07_arch(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide)
    add_title(slide, "Arquitectura técnica", "Stack de Ingeniería de Datos completo y trazable")

    add_image_safe(slide, "architecture.png",
                   Inches(0.8), Inches(1.7), width=Inches(8.0))

    # right column: tech stack
    tx = Inches(9.2); ty = Inches(1.85)
    add_text(slide, "STACK", tx, ty, Inches(3.5), Inches(0.3),
             size=10, bold=True, color=CORAL)
    techs = [
        ("Extract", "Socrata SODA + Parquet", NAVY),
        ("Store", "PostgreSQL + MongoDB", NAVY),
        ("Compute", "Python · uv · Polars", TEAL),
        ("Serve", "FastAPI (3 servicios)", TEAL),
        ("UI", "Dash + Plotly", CORAL),
        ("Orquesta", "Make + Docker Compose", CORAL),
    ]
    yy = ty + Inches(0.35)
    for label, val, color in techs:
        add_text(slide, label, tx, yy, Inches(1.2), Inches(0.3),
                 size=10, bold=True, color=color)
        add_text(slide, val, tx + Inches(1.3), yy, Inches(2.4), Inches(0.3),
                 size=11, color=INK)
        yy += Inches(0.42)

    add_rect(slide, Inches(0.7), Inches(6.5), Inches(12), Inches(0.5), SAND)
    add_text(slide,
             "PostgreSQL es fuente de verdad relacional. MongoDB guarda evidencia documental y eventos.",
             Inches(0.85), Inches(6.6), Inches(11.7), Inches(0.4),
             size=11, italic=True, color=INK)

    add_footer(slide, 7, TOTAL)


def s08_score(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide)
    add_title(slide, "Score explicable", "Cuatro componentes auditables → un número entre 0 y 100")

    components = [
        ("Reglas",     "Heurísticas verificables\n(modalidad, unicidad, plazos).", NAVY),
        ("Pares",      "Desviación contra procesos\ncomparables del segmento.",   TEAL),
        ("Anomalía",   "Componente estadístico\n(Isolation Forest).",              CORAL),
        ("Confianza",  "0–1 según cobertura\ny calidad de datos.",                MUTED),
    ]
    cw = Inches(2.9); ch = Inches(2.4)
    sx = Inches(0.7); gap = Inches(0.2)
    for i, (head, body, color) in enumerate(components):
        x = sx + i * (cw + gap)
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x, Inches(1.9), cw, ch)
        shp.fill.solid(); shp.fill.fore_color.rgb = WHITE
        shp.line.color.rgb = MIST; shp.line.width = Pt(0.75); shp.adjustments[0] = 0.05
        add_rect(slide, x, Inches(1.9), cw, Pt(4), color)
        add_text(slide, head, x + Inches(0.25), Inches(2.05),
                 cw - Inches(0.5), Inches(0.5), size=15, bold=True, color=INK)
        add_text(slide, body, x + Inches(0.25), Inches(2.65),
                 cw - Inches(0.5), Inches(1.7), size=11, color=MUTED)

    # interpretation strip
    add_rect(slide, Inches(0.7), Inches(4.7), Inches(12), Inches(2.0), CLOUD)
    add_text(slide, "INTERPRETACIÓN DEL SCORE",
             Inches(0.95), Inches(4.85), Inches(11), Inches(0.3),
             size=10, bold=True, color=CORAL)

    ranges = [
        ("0–20",  "Típico",          "Sin desviaciones relevantes."),
        ("21–40", "Leve",            "Revisar solo si hay contexto sensible."),
        ("41–70", "Notable",         "Revisión recomendada con evidencia."),
        ("71–100","Alta prioridad",  "Multi-señal, revisar primero."),
    ]
    rw = Inches(2.85); rh = Inches(1.3)
    ry = Inches(5.25)
    for i, (rng, lbl, body) in enumerate(ranges):
        x = Inches(0.85) + i * (rw + Inches(0.1))
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, ry, rw, rh)
        shp.fill.solid(); shp.fill.fore_color.rgb = WHITE
        shp.line.fill.background(); shp.adjustments[0] = 0.08
        color = [TEAL, NAVY, RGBColor(0xC2,0x84,0x2A), CORAL][i]
        add_text(slide, rng, x, ry + Inches(0.1), rw, Inches(0.35),
                 size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, lbl, x, ry + Inches(0.45), rw, Inches(0.3),
                 size=10, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + Inches(0.1), ry + Inches(0.75), rw - Inches(0.2), Inches(0.5),
                 size=9, color=MUTED, align=PP_ALIGN.CENTER)

    add_footer(slide, 8, TOTAL)


def s09_engineering(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide)
    add_title(slide, "Evidencia de ingeniería", "Métricas duras de un sistema reproducible")

    # left: big metrics grid
    metrics = [
        ("17.229", "procesos cargados", NAVY),
        ("33",     "objetos PostgreSQL", TEAL),
        ("5/5",    "colecciones MongoDB", CORAL),
        ("3",      "APIs FastAPI · health 200", NAVY),
        ("71",     "tests pytest", TEAL),
        ("21",     "documentos técnicos", CORAL),
    ]
    cols = 3; rows = 2
    mw = Inches(2.45); mh = Inches(1.55)
    gx = Inches(0.15); gy = Inches(0.15)
    sx = Inches(0.7); sy = Inches(1.95)
    for i, (v, l, c) in enumerate(metrics):
        r = i // cols; col = i % cols
        x = sx + col * (mw + gx)
        y = sy + r * (mh + gy)
        add_metric(slide, x, y, mw, mh, v, l, value_color=c, bg=WHITE, value_size=28)
        # add subtle border
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, mw, mh)
        shp.fill.background()
        shp.line.color.rgb = MIST; shp.line.width = Pt(0.5); shp.adjustments[0] = 0.08

    # right: SQL engineering call-outs
    rx = Inches(8.7); ry = Inches(1.95)
    add_text(slide, "SQL ENGINEERING VISIBLE",
             rx, ry, Inches(4.5), Inches(0.3),
             size=10, bold=True, color=CORAL)
    sql_items = [
        ("Triggers",         "audit_log · historial · updated_at"),
        ("Window functions", "concentración · outliers"),
        ("CTE recursiva",    "jerarquía territorial"),
        ("Transacciones",    "score + evento atómico"),
        ("Vistas",           "ranking + alineación PAA"),
    ]
    yy = ry + Inches(0.4)
    for label, val in sql_items:
        add_text(slide, label, rx, yy, Inches(1.4), Inches(0.3),
                 size=10, bold=True, color=NAVY)
        add_text(slide, val, rx + Inches(1.5), yy, Inches(3.3), Inches(0.3),
                 size=10, color=MUTED)
        yy += Inches(0.45)

    add_footer(slide, 9, TOTAL)


def s10_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide)
    add_title(slide, "Demo guiada", "Tres clics: panorama → ranking → evidencia")

    shots = [
        ("screenshot_dashboard_home.png", "1", "Panorama",
         "El usuario entiende\nel universo a revisar."),
        ("screenshot_ranking.png", "2", "Ranking",
         "Procesos ordenados por\nprioridad, con confianza."),
        ("screenshot_process_detail.png", "3", "Detalle",
         "Razones, comparables\ny reporte exportable."),
    ]
    iw = Inches(3.85); ih = Inches(2.7)
    sx = Inches(0.7); gap = Inches(0.3); sy = Inches(1.95)
    for i, (fname, num, head, body) in enumerate(shots):
        x = sx + i * (iw + gap)
        # screenshot frame
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, sy, iw, ih)
        shp.fill.solid(); shp.fill.fore_color.rgb = WHITE
        shp.line.color.rgb = MIST; shp.line.width = Pt(0.75); shp.adjustments[0] = 0.03
        add_image_safe(slide, fname, x + Inches(0.1), sy + Inches(0.1),
                       width=iw - Inches(0.2), height=ih - Inches(0.2))
        # caption block
        cy = sy + ih + Inches(0.15)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, cy, Inches(0.45), Inches(0.45))
        circle.fill.solid(); circle.fill.fore_color.rgb = CORAL
        circle.line.fill.background()
        add_text(slide, num, x, cy, Inches(0.45), Inches(0.45),
                 size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, head, x + Inches(0.55), cy + Inches(0.02),
                 iw - Inches(0.6), Inches(0.4), size=14, bold=True, color=INK)
        add_text(slide, body, x + Inches(0.55), cy + Inches(0.45),
                 iw - Inches(0.6), Inches(1.0), size=10, color=MUTED)

    add_footer(slide, 10, TOTAL)


def s11_validation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide)
    add_title(slide, "Validación y límites", "Qué validamos, qué falta, próximos pasos honestos")

    # validation chart
    add_image_safe(slide, "validation_summary.png",
                   Inches(0.7), Inches(1.95), width=Inches(5.6))

    # right side
    rx = Inches(6.8); ry = Inches(1.95)
    add_text(slide, "VALIDACIÓN FINAL", rx, ry, Inches(6), Inches(0.3),
             size=10, bold=True, color=TEAL)
    add_text(slide, "make validate-final  ·  ok",
             rx, ry + Inches(0.3), Inches(6), Inches(0.5),
             size=18, bold=True, color=TEAL, font="Consolas")

    items = [
        ("Software", "71 tests pytest, lint estable, validación reproducible."),
        ("Datos",    "17.229 procesos cargados, joins auditados, fill rates documentados."),
        ("Compuerta PAA", "Match top-1 con confianza ≥ 0.75 reportado por entidad y modalidad."),
    ]
    yy = ry + Inches(1.1)
    for head, body in items:
        add_text(slide, "✓  " + head, rx, yy, Inches(6), Inches(0.3),
                 size=12, bold=True, color=INK)
        add_text(slide, body, rx + Inches(0.3), yy + Inches(0.3),
                 Inches(5.8), Inches(0.5), size=10, color=MUTED)
        yy += Inches(0.75)

    # honest limits strip
    add_rect(slide, Inches(0.7), Inches(6.0), Inches(12), Inches(0.95), SAND)
    add_text(slide, "LÍMITES Y SIGUIENTE PASO",
             Inches(0.9), Inches(6.1), Inches(11), Inches(0.3),
             size=10, bold=True, color=CORAL)
    add_text(slide,
             "Validación humana de 100 casos + 5 entrevistas reales + backend de embeddings activable: "
             "ya hay protocolo, faltan revisores. Honesto > teatro.",
             Inches(0.9), Inches(6.4), Inches(11.6), Inches(0.5),
             size=11, color=INK)

    add_footer(slide, 11, TOTAL)


def s12_impact(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide)
    add_title(slide, "Impacto y escalabilidad", "De la Orinoquía a Colombia, con ética por diseño")

    # left: impact narrative
    add_rect(slide, Inches(0.7), Inches(1.9), Inches(0.18), Inches(4.5), CORAL)
    add_text(slide, "DEMO ACTUAL", Inches(1.0), Inches(1.9),
             Inches(5), Inches(0.3), size=10, bold=True, color=CORAL)
    add_text(slide,
             "Meta y Casanare — Orinoquía como banco de pruebas territorial: "
             "una oficina de control con capacidad limitada convierte 17.229 procesos "
             "en una cola priorizada de 50 casos revisables con razones visibles.",
             Inches(1.0), Inches(2.25), Inches(5.5), Inches(2.0),
             size=12, color=INK)

    add_text(slide, "PRÓXIMOS PASOS", Inches(1.0), Inches(4.5),
             Inches(5), Inches(0.3), size=10, bold=True, color=CORAL)
    steps = [
        "Pilotos con 2 oficinas de control en Orinoquía.",
        "Activar backend de embeddings (sentence-transformers).",
        "Despliegue público read-only (Render/Fly).",
        "Escalado por departamento, mismo pipeline.",
    ]
    yy = Inches(4.85)
    for s in steps:
        add_text(slide, "→  " + s, Inches(1.0), yy, Inches(5.5), Inches(0.32),
                 size=11, color=INK)
        yy += Inches(0.38)

    # right: ethics by design
    rx = Inches(7.2); ry = Inches(1.9)
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 rx, ry, Inches(5.5), Inches(4.8))
    shp.fill.solid(); shp.fill.fore_color.rgb = INK
    shp.line.fill.background(); shp.adjustments[0] = 0.04

    add_text(slide, "ÉTICA POR DISEÑO", rx + Inches(0.3), ry + Inches(0.25),
             Inches(5), Inches(0.3), size=10, bold=True, color=CORAL)
    add_text(slide, "Cinco compromisos no negociables",
             rx + Inches(0.3), ry + Inches(0.6),
             Inches(5), Inches(0.5), size=16, bold=True, color=WHITE)

    ethics = [
        "Disclaimer visible en UI, reportes y exportes.",
        "Revisión humana obligatoria antes de cualquier acción.",
        "Trazabilidad: fuente, fecha, hash, snapshot.",
        "Contexto fiscal NO entra como etiqueta del modelo.",
        "Modelo y datos auditables: pesos y reglas públicos.",
    ]
    yy = ry + Inches(1.55)
    for it in ethics:
        add_text(slide, "•  " + it, rx + Inches(0.3), yy,
                 Inches(5), Inches(0.45), size=11, color=CLOUD)
        yy += Inches(0.55)

    add_footer(slide, 12, TOTAL)


def s13_close(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, INK)
    add_rect(slide, 0, 0, Inches(0.25), SLIDE_H, CORAL)

    add_text(slide, "ContratIA Abierta",
             Inches(0.9), Inches(1.0), Inches(11), Inches(1.2),
             size=48, bold=True, color=WHITE)

    add_text(slide,
             "Prioriza revisión humana en la contratación pública colombiana.",
             Inches(0.9), Inches(2.4), Inches(11), Inches(0.6),
             size=18, color=CLOUD)

    add_text(slide,
             "No emitimos juicios de corrupción. Ordenamos qué revisar primero, con evidencia trazable.",
             Inches(0.9), Inches(3.05), Inches(11), Inches(0.5),
             size=13, italic=True, color=SAND)

    add_rect(slide, Inches(0.9), Inches(3.75), Inches(1.2), Pt(2), CORAL)

    # info grid
    add_text(slide, "EQUIPO", Inches(0.9), Inches(4.1),
             Inches(4), Inches(0.3), size=10, bold=True, color=CORAL)
    add_text(slide, "Ingeniería de Datos · Universidad del Rosario",
             Inches(0.9), Inches(4.4), Inches(7), Inches(0.4),
             size=14, color=WHITE)

    add_text(slide, "REPOSITORIO", Inches(0.9), Inches(5.0),
             Inches(4), Inches(0.3), size=10, bold=True, color=CORAL)
    add_text(slide, "github.com/Thom-320/secop-risk-alerts-co",
             Inches(0.9), Inches(5.3), Inches(7), Inches(0.4),
             size=14, color=WHITE, font="Consolas")

    add_text(slide, "CATEGORÍA  ·  REGIÓN DEMO", Inches(0.9), Inches(5.9),
             Inches(6), Inches(0.3), size=10, bold=True, color=CORAL)
    add_text(slide, "Gobernanza y Transparencia  ·  Meta y Casanare",
             Inches(0.9), Inches(6.2), Inches(8), Inches(0.4),
             size=14, color=WHITE)

    # right hero quote
    add_rect(slide, Inches(8.6), Inches(4.1), Inches(4.1), Inches(2.6), NAVY)
    add_text(slide,
             "“Convierte\nuna revisión\nimposible\nen una cola\nauditable.”",
             Inches(8.8), Inches(4.25), Inches(3.8), Inches(2.4),
             size=18, bold=True, color=WHITE, italic=True)


# ── Entry ──────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    s01_cover(prs)
    s02_problem(prs)
    s03_thesis(prs)
    s04_user(prs)
    s05_data(prs)
    s06_how(prs)
    s07_arch(prs)
    s08_score(prs)
    s09_engineering(prs)
    s10_demo(prs)
    s11_validation(prs)
    s12_impact(prs)
    s13_close(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"OK — {len(prs.slides)} slides → {OUTPUT}")


if __name__ == "__main__":
    main()
