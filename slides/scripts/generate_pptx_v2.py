#!/usr/bin/env python3
"""Generate ContratIA Abierta / Transparencia360 PPTX deck v2."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(0xF8, 0xFA, 0xFC)
ACCENT = RGBColor(0x27, 0x5D, 0xA8)
DARK = RGBColor(0x26, 0x32, 0x41)
MUTED = RGBColor(0x66, 0x70, 0x85)
RULE = RGBColor(0xD9, 0xDE, 0xE8)
SOFT_BG = RGBColor(0xEA, 0xF1, 0xFB)
GREEN = RGBColor(0x2F, 0x8F, 0x64)
AMBER = RGBColor(0xB7, 0x79, 0x1F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"
ASSETS = Path(__file__).resolve().parent.parent / "assets"
OUTPUT = Path(__file__).resolve().parent.parent / "contratia_abierta_deck_v2.pptx"

FOOTER_TEXT = "Priorización de revisión humana; no prueba conductas indebidas."


def set_slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_bar(slide):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Pt(4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_slide_number(slide, num):
    txBox = slide.shapes.add_textbox(
        SLIDE_W - Inches(0.8), SLIDE_H - Inches(0.45), Inches(0.6), Inches(0.3)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED
    p.font.name = FONT
    p.alignment = PP_ALIGN.RIGHT


def add_footer(slide):
    txBox = slide.shapes.add_textbox(
        Inches(0.6), SLIDE_H - Inches(0.45), Inches(8), Inches(0.3)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = FOOTER_TEXT
    p.font.size = Pt(8)
    p.font.color.rgb = MUTED
    p.font.name = FONT
    p.font.italic = True


def add_title(slide, text, top=Inches(0.5), left=Inches(0.6), width=Inches(12), size=Pt(28)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = True
    p.font.color.rgb = DARK
    p.font.name = FONT
    return txBox


def add_text(slide, text, left, top, width, height, size=Pt(14), color=DARK,
             bold=False, italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.name = FONT
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = align
    try:
        tf.paragraphs[0].space_after = Pt(0)
        tf.paragraphs[0].space_before = Pt(0)
    except Exception:
        pass
    return txBox


def add_rich_text(slide, runs_list, left, top, width, height, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for i, run_spec in enumerate(runs_list):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
        else:
            run = p.add_run()
        run.text = run_spec.get("text", "")
        run.font.size = run_spec.get("size", Pt(14))
        run.font.color.rgb = run_spec.get("color", DARK)
        run.font.name = FONT
        run.font.bold = run_spec.get("bold", False)
        run.font.italic = run_spec.get("italic", False)
    return txBox


def add_card(slide, left, top, width, height, title_text, body_text="",
             title_color=DARK, body_color=MUTED, bg_color=WHITE, border_color=RULE,
             title_size=Pt(13), body_size=Pt(11)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    shape.adjustments[0] = 0.04

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)

    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = title_size
    p.font.bold = True
    p.font.color.rgb = title_color
    p.font.name = FONT
    p.space_after = Pt(4)

    if body_text:
        p2 = tf.add_paragraph()
        p2.text = body_text
        p2.font.size = body_size
        p2.font.color.rgb = body_color
        p2.font.name = FONT
        p2.space_before = Pt(2)

    return shape


def add_metric_box(slide, left, top, width, height, value, label,
                   value_color=ACCENT, bg_color=SOFT_BG):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    shape.adjustments[0] = 0.06

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)

    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = value_color
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(9)
    p2.font.color.rgb = MUTED
    p2.font.name = FONT
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(2)

    return shape


def add_chip(slide, left, top, text, bg=SOFT_BG, fg=ACCENT):
    w = Inches(max(1.2, len(text) * 0.1 + 0.3))
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    shape.adjustments[0] = 0.5

    tf = shape.text_frame
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = fg
    p.font.name = FONT
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_arrow(slide, left, top, width=Inches(0.4)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.25)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


def add_image_safe(slide, name, left, top, width=None, height=None):
    path = ASSETS / name
    if not path.exists():
        add_text(slide, f"[{name}]", left, top,
                 width or Inches(4), height or Inches(2),
                 size=Pt(11), color=MUTED, italic=True)
        return None
    kwargs = {}
    if width:
        kwargs["width"] = width
    if height:
        kwargs["height"] = height
    return slide.shapes.add_picture(str(path), left, top, **kwargs)


def add_table(slide, left, top, width, height, headers, rows,
              header_bg=ACCENT, header_fg=WHITE, row_bg=WHITE, alt_bg=SOFT_BG):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    col_w = int(width / n_cols)
    for i in range(n_cols):
        table.columns[i].width = col_w

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = header_fg
            p.font.name = FONT
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg

    for r_idx, row in enumerate(rows):
        bg = alt_bg if r_idx % 2 == 1 else row_bg
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = DARK
                p.font.name = FONT
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

    return table_shape


def add_rule_line(slide, left, top, width, color=RULE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ── Slide builders ──────────────────────────────────────────────

def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide)

    add_text(slide, "Transparencia360", Inches(0.7), Inches(1.0), Inches(7), Inches(0.8),
             size=Pt(42), bold=True, color=DARK)
    add_text(slide, "ContratIA Abierta", Inches(0.7), Inches(1.8), Inches(7), Inches(0.7),
             size=Pt(34), bold=True, color=DARK)
    add_text(slide, "Ordenar miles de procesos SECOP para decidir qué revisar primero.",
             Inches(0.7), Inches(2.8), Inches(7), Inches(0.5),
             size=Pt(16), bold=True, color=ACCENT)

    add_rule_line(slide, Inches(0.7), Inches(3.5), Inches(5))

    add_text(slide, "Ingeniería de Datos, Universidad del Rosario",
             Inches(0.7), Inches(3.7), Inches(6), Inches(0.4),
             size=Pt(12), color=MUTED)

    chip_x = Inches(0.7)
    for tag in ["Revisión humana", "Datos abiertos", "Trazabilidad"]:
        c = add_chip(slide, chip_x, Inches(4.3), tag)
        chip_x += c.width + Inches(0.15)

    add_image_safe(slide, "architecture.png",
                   Inches(8.5), Inches(1.5), width=Inches(4.2))

    add_slide_number(slide, 1)


def slide_02_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide)
    add_title(slide, "Problema: volumen sin atención suficiente")

    add_text(slide,
             "No faltan datos de contratación. Falta capacidad humana para revisarlos "
             "con prioridad, evidencia y trazabilidad.",
             Inches(0.6), Inches(1.5), Inches(12), Inches(0.6),
             size=Pt(15), color=DARK)

    box_w = Inches(3.4)
    box_h = Inches(2.6)
    gap = Inches(0.5)
    start_x = Inches(0.8)
    box_y = Inches(2.8)

    boxes = [
        ("Miles de procesos", "SECOP, PAA, contexto"),
        ("Ranking explicable", "Score, confianza, razones"),
        ("Revisión humana", "Decisión documentada"),
    ]

    for i, (title, body) in enumerate(boxes):
        x = start_x + i * (box_w + gap + Inches(0.4))
        add_card(slide, x, box_y, box_w, box_h, title, body,
                 title_color=ACCENT, bg_color=WHITE)
        if i < 2:
            arrow_x = x + box_w + Inches(0.05)
            add_arrow(slide, arrow_x, box_y + Inches(1.1))

    add_footer(slide)
    add_slide_number(slide, 2)


def slide_03_stakeholders(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide)
    add_title(slide, "Stakeholder y decisión")

    cards = [
        ("Veeduría ciudadana", "Monitorea contratos públicos y exige transparencia."),
        ("Control interno", "Audita procesos y prioriza revisiones internas."),
        ("Periodista de datos", "Investiga patrones y publica hallazgos."),
    ]

    card_w = Inches(3.6)
    card_h = Inches(2.4)
    start_x = Inches(0.8)
    gap = Inches(0.5)

    for i, (t, b) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_card(slide, x, Inches(1.8), card_w, card_h, t, b,
                 title_color=ACCENT, bg_color=SOFT_BG)

    add_rule_line(slide, Inches(0.6), Inches(5.0), Inches(12))

    add_text(slide,
             "Decisión semanal: qué procesos revisar primero, con qué razones y con qué confianza.",
             Inches(0.6), Inches(5.3), Inches(12), Inches(0.5),
             size=Pt(14), bold=True, color=ACCENT)

    add_footer(slide)
    add_slide_number(slide, 3)


def slide_04_requirements(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide)
    add_title(slide, "Requisitos de Ingeniería de Datos cubiertos")

    items = [
        ("PostgreSQL", "27 tablas, 33 objetos"),
        ("MongoDB", "5 colecciones"),
        ("FastAPI", "3 servicios, health 200"),
        ("Dash", "Interfaz oficial"),
        ("ETL", "17.229 procesos"),
        ("SQL avanzado", "Triggers, CTE, windows"),
        ("Pruebas", "66 pytest pasan"),
        ("Documentación", "21 archivos"),
    ]

    cols = 4
    rows = 2
    card_w = Inches(2.7)
    card_h = Inches(1.8)
    gap_x = Inches(0.35)
    gap_y = Inches(0.35)
    start_x = Inches(0.6)
    start_y = Inches(1.7)

    for idx, (t, b) in enumerate(items):
        r = idx // cols
        c = idx % cols
        x = start_x + c * (card_w + gap_x)
        y = start_y + r * (card_h + gap_y)
        add_card(slide, x, y, card_w, card_h, t, b,
                 title_color=ACCENT, bg_color=WHITE, title_size=Pt(12), body_size=Pt(10))

    add_footer(slide)
    add_slide_number(slide, 4)


def slide_05_data_sources(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide)
    add_title(slide, "Datos oficiales y volumen de demo")

    headers = ["Dataset", "Fuente", "Procesos", "Uso"]
    rows = [
        ["p6dx", "SECOP II", "12.406", "Contratos"],
        ["rpmr", "SECOP Integrado", "4.823", "Histórico"],
        ["9sue", "PAA", "4.494 items", "Planeación"],
        ["wasc", "Control fiscal", "—", "Auditoría"],
    ]

    add_table(slide, Inches(0.6), Inches(1.6), Inches(6.5), Inches(3.2),
              headers, rows)

    metrics = [
        ("17.229", "Procesos"),
        ("68.916", "Razones"),
        ("4.494", "PAA items"),
        ("41.996", "Eventos"),
    ]

    mx = Inches(7.8)
    my = Inches(1.6)
    mw = Inches(2.4)
    mh = Inches(1.3)
    gap = Inches(0.25)

    for i, (val, label) in enumerate(metrics):
        r = i // 2
        c = i % 2
        x = mx + c * (mw + gap)
        y = my + r * (mh + gap)
        add_metric_box(slide, x, y, mw, mh, val, label)

    add_footer(slide)
    add_slide_number(slide, 5)


def slide_06_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide)
    add_title(slide, "Arquitectura de extremo a extremo")

    add_image_safe(slide, "architecture.png",
                   Inches(1.0), Inches(1.6), width=Inches(11.3))

    add_text(slide, "Socrata → ETL → PostgreSQL + MongoDB → FastAPI → Dash",
             Inches(0.6), Inches(6.2), Inches(12), Inches(0.4),
             size=Pt(11), color=MUTED, italic=True, align=PP_ALIGN.CENTER)

    add_footer(slide)
    add_slide_number(slide, 6)


def slide_07_er_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide)
    add_title(slide, "Modelo relacional como fuente de verdad")

    add_image_safe(slide, "er_model.png",
                   Inches(0.8), Inches(1.5), width=Inches(11.7))

    add_text(slide,
             "27 tablas relacionales con PK/FK, constraints e índices; 33 objetos públicos incluyendo vistas.",
             Inches(0.6), Inches(6.2), Inches(12), Inches(0.4),
             size=Pt(11), color=MUTED, italic=True, align=PP_ALIGN.CENTER)

    add_footer(slide)
    add_slide_number(slide, 7)


def slide_08_nosql(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide)
    add_title(slide, "NoSQL y auditoría documental")

    headers = ["Colección", "Documentos", "Propósito"]
    rows = [
        ["raw_process_snapshots", "100", "Snapshot crudo"],
        ["etl_run_logs", "1", "Log de ejecución ETL"],
        ["risk_event_logs", "1", "Eventos de riesgo"],
        ["report_snapshots", "1", "Reportes generados"],
        ["user_action_logs", "1", "Acciones de usuario"],
    ]

    add_table(slide, Inches(0.6), Inches(1.6), Inches(6.5), Inches(3.8),
              headers, rows)

    add_text(slide,
             "MongoDB guarda evidencia flexible\nsin deformar el modelo relacional.",
             Inches(7.8), Inches(2.5), Inches(4.5), Inches(1.5),
             size=Pt(15), color=DARK, bold=True)

    add_footer(slide)
    add_slide_number(slide, 8)


def slide_09_sql(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide)
    add_title(slide, "SQL engineering visible")

    items = [
        ("Triggers", "audit_log, historial, updated_at"),
        ("Window functions", "Concentración y outliers"),
        ("CTE recursiva", "Jerarquía territorial"),
        ("Transacciones", "Score + evento atómico"),
    ]

    card_w = Inches(5.5)
    card_h = Inches(1.8)
    gap_x = Inches(0.5)
    gap_y = Inches(0.4)
    start_x = Inches(0.8)
    start_y = Inches(1.7)

    for idx, (t, b) in enumerate(items):
        r = idx // 2
        c = idx % 2
        x = start_x + c * (card_w + gap_x)
        y = start_y + r * (card_h + gap_y)
        add_card(slide, x, y, card_w, card_h, t, b,
                 title_color=ACCENT, bg_color=WHITE, title_size=Pt(14), body_size=Pt(12))

    add_rule_line(slide, Inches(0.6), Inches(5.8), Inches(12))

    add_text(slide,
             "No es solo almacenamiento: hay integridad, trazabilidad y consultas analíticas reproducibles.",
             Inches(0.6), Inches(6.0), Inches(12), Inches(0.4),
             size=Pt(12), bold=True, color=DARK)

    add_footer(slide)
    add_slide_number(slide, 9)


def slide_10_score(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide)
    add_title(slide, "Score explicable, no veredicto")

    cards = [
        ("Reglas", "Heurísticas de\nriesgo contractivo"),
        ("Pares", "Comparables por\nsegmento y valor"),
        ("Anomalía", "Desviación estadística\ndel patrón"),
        ("Confianza", "0–1 según cobertura\nde datos"),
        ("Razones", "Texto legible por\nhumano"),
    ]

    card_w = Inches(2.2)
    card_h = Inches(2.6)
    gap = Inches(0.25)
    total_w = 5 * card_w + 4 * gap
    start_x = (SLIDE_W - total_w) // 2

    for i, (t, b) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_card(slide, x, Inches(1.8), card_w, card_h, t, b,
                 title_color=ACCENT, bg_color=SOFT_BG, title_size=Pt(13), body_size=Pt(10))

    add_rule_line(slide, Inches(0.6), Inches(5.2), Inches(12))

    add_text(slide,
             "Salida: score de prioridad + confianza + razones + comparables.",
             Inches(0.6), Inches(5.5), Inches(12), Inches(0.4),
             size=Pt(14), bold=True, color=ACCENT)

    add_footer(slide)
    add_slide_number(slide, 10)


def slide_11_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide)
    add_title(slide, "Demo: flujo de revisión")

    screenshots = [
        ("screenshot_dashboard_home.png", "1  Panorama"),
        ("screenshot_ranking.png", "2  Ranking"),
        ("screenshot_process_detail.png", "3  Detalle"),
    ]

    img_w = Inches(3.6)
    img_h = Inches(3.6)
    gap = Inches(0.4)
    total_w = 3 * img_w + 2 * gap
    start_x = (SLIDE_W - total_w) // 2

    for i, (fname, label) in enumerate(screenshots):
        x = start_x + i * (img_w + gap)
        add_image_safe(slide, fname, x, Inches(1.6), width=img_w, height=img_h)
        add_text(slide, label, x, Inches(5.4), img_w, Inches(0.35),
                 size=Pt(11), color=MUTED, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide,
             "El usuario entiende el universo, filtra candidatos y revisa evidencia trazable.",
             Inches(0.6), Inches(6.0), Inches(12), Inches(0.4),
             size=Pt(12), color=DARK, italic=True, align=PP_ALIGN.CENTER)

    add_footer(slide)
    add_slide_number(slide, 11)


def slide_12_validation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide)
    add_title(slide, "Validación, límites y siguiente paso")

    add_image_safe(slide, "validation_summary.png",
                   Inches(0.6), Inches(1.6), width=Inches(5.8))

    add_text(slide, "Validación final: ok",
             Inches(7.2), Inches(1.8), Inches(5.5), Inches(0.5),
             size=Pt(20), bold=True, color=GREEN)

    add_rule_line(slide, Inches(7.2), Inches(2.5), Inches(5))

    add_text(slide,
             "Límites: datos ruidosos, joins imperfectos,\nrequiere revisión humana.",
             Inches(7.2), Inches(2.8), Inches(5.5), Inches(1.2),
             size=Pt(13), color=DARK)

    add_text(slide,
             "Siguiente: encuesta real con 5 usuarios\ny piloto controlado.",
             Inches(7.2), Inches(4.2), Inches(5.5), Inches(1.0),
             size=Pt(14), bold=True, color=ACCENT)

    add_footer(slide)
    add_slide_number(slide, 12)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_stakeholders(prs)
    slide_04_requirements(prs)
    slide_05_data_sources(prs)
    slide_06_architecture(prs)
    slide_07_er_model(prs)
    slide_08_nosql(prs)
    slide_09_sql(prs)
    slide_10_score(prs)
    slide_11_demo(prs)
    slide_12_validation(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"OK — {len(prs.slides)} slides written to {OUTPUT}")


if __name__ == "__main__":
    main()
